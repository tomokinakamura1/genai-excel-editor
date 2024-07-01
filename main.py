from pptx import Presentation 

path = '/Users/karthik/Documents/GitHub/hipem/ppt/HONDA_FTA_1pager.pptx'
text_content  = ''
ppt = Presentation(path)

for slide_number, slide in enumerate(ppt.slides): 
    for shape in slide.shapes: 
        if hasattr(shape, "text"): 
            text_content += shape.text + "\n"
            print(text_content)

# with open('extracted_content.txt', "w") as f:
#     f.write(text_content)