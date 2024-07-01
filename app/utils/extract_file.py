from ibm_watson_machine_learning.foundation_models import Model
import pandas as pd
import json
import os
from pptx import Presentation
from dotenv import load_dotenv
from app.prompts import v1

load_dotenv()


stars_data = pd.read_excel("tables.xlsx")


def get_credentials():
    return {"url": os.getenv("IBM_CLOUD_URL"), "apikey": os.getenv("IBM_CLOUD_API_KEY")}


parameters = {
    "decoding_method": "greedy",
    "max_new_tokens": 2000,
    "min_new_tokens": 1,
    "stop_sequences": ["<end>"],
}


def create_model():
    model = Model(
        model_id="meta-llama/llama-3-70b-instruct",
        params=parameters,
        credentials=get_credentials(),
        project_id=os.getenv("PROJECT_ID"),
    )
    return model


model = create_model()

input_metadata_question = {
    "Customer Problem:": "They have 10 years of fault trouble analysis (FTA) tree to analyze the root cause of car trouble. But, those tree are contained in Power Point. So, it’s hard for junior engineer to find the appropriate document and list up all possible root cause from the tree. \n\nWe proposed the extraction method of tree in original algorithm and to use watsonx.data to search the tree database efficiently.",
    "IBM Portfolio Products Used:": "watsonx.data\nwatsonx Orchestrate",
    "IBM Team Engaged:": "Market: Japan\nCE Leader: Toru Nakayama (FLM)\nWorkshop: Not conducted\nCE AI Engineers: Ryohei Fujii\nIBM Consulting Included: Yes. The client’s department is quality-related one. IBM-C engaged in the other quality AI system in the past.",
    "Outcome:": "The client struggled to digitalize their existing tree data in PowerPoint, and we solve this problem in this Pilot. They’ve highly satisfied with our proposal because their FTA assets have been stored for 10 years. FTA information is very important for analyzing the root cause of car trouble, so this is key factor to start the digital transformation of their quality area. \nCurrently, they took 1,200 days from they find the issue and until they get the action, and cars which has a potential trouble continue to be manufactured. They aims to reduce this time to 600 days. The economic impact is very huge if we can reduce the amount of potential-risk car for the rest of 600 days.",
    "Next Steps:": "They’ve satisfied with our result, and they talked about their all strategy of quality DX. In addition, they planned to start PoC from this July, and asked us to propose our plan which collaborate with the client. Currently, CE and IBM-C team consider the next PoC plan.",
}

internal_pre_sales_data_info = {
    "Project Description(limit to 270 Characters)": "To extract the expert engineer’s knowledge to improve the A-ES in HONDA. A-ES is the knowledge retrieval system based on graph network but it takes a long time to expand the other area. We proposed the generative AI platform to solve this problem",
    "Engagement and Use Case Summary(limit to 2000 Characters)": "To extract the expert engineer’s knowledge to improve the A-ES in HONDA. A-ES is the knowledge retrieval system based on graph network but it takes a long time to expand the other area. We proposed the generative AI platform to solve this problem.",
    "Business Value and Proposition, ROI, Competitive Differentiators (limit to 2000 Characters)": """We have built the secure environment for treating client’s confidential dataset in IBM Cloud. This is a first time for this client, and our client can get the environment for testing the new technology with IBM teams (CE/IBM-C/TEL/CSM).This environment accelerate the time to solve their business concern.And, we define the framework for testing their document based on solid criteria, and this framework helps to continue and improve this activity.In the behalf of the project, we tested  our proposed method by using their confidential dataset, and our client know the best use case and bad use case for the generative AI. """,
    "Stories / Narratives(limit to 2000 Characters)": "The client’s key person have been led A-ES project, and wanted to expand the area of A-ES widely. Current system has required to take a lot of time to create the model which can fit to the existing system. We heard that the client’s experts remains their own knowledge in the Power Point decks, but there are many insufficient text and many unstructured components like tables and images. So, we proposed the new method to extract useful informations from such documents. If we have more plenty of resource about GPU on SaaS or On-premiss in Techzone, we can realize more difficult task and satisfy the client’s needs. ",
    "Engagement Success(limit to 2000 Characters)": "Our initial goal is to verify the possibility of extracting the knowledge from the unstructured components like graphs and images. We tried to use the multi-modal model and show the summary result of each page to the client. As a result, the client felt the feasibility of this concept and thought that our proposed new approach will jump up to the next step of the roadmap of their knowledge system.  ",
    "Top Lessons Learned and Failing Forward(limit to 1000 Characters)": """Good point: The challenging proposal leads to high satisfaction of the client. The client decided to move to the paid project by looking our result and get trust to IBM team. Client Engineering is no-charge team, so we can challenge the difficult and new approach. Bad point: We took a month to apply the custom request of Techzone to treat the customer’s secure dataset. Unfortunately, we took longer a month to apply the client-side agreement of using their secure dataset in the custom environment due to the content of Client Engineering Engagement which said “We charge only $1000 if the incident happened“ and the ownership of the testing environment is IBM (Techzone). So, we learned that the data owner is different of LOB users and CE activity is not fitted to the secure data use case.  """,
    "Client Feedback(limit to 1000 Characters)": """We got the client comment in client story of this project. (This will be published soon.)" IBM's innovative solutions and the platform that securely leverages our vast amount of development information are contributing to our dream of delivering more value to our customers. We look forward to further cooperation. """,
}
close_out_transition_info = {
    "Client Technology Landscape(limit to 2000 Characters)": "Our client technology landscape related with the generative AI is very high. They’ve prepared Azure OpenAI service environment on their secure Azure Cloud and enterprise-level ChatGPT service in the early phase. But they didn’t try image-text model due to the lack of the capability of the above environment, so we tried it and said that “watsonx.ai has model-free environment and can choose the appropriate model for your use case from OS model, 3rd party model, and IBM model”. In addition, the speed and technology level of CE and the business knowledge of IBM-C who have engaged with clients in long years are recognized. ",
    "Next Steps for IBM(limit to 1000 Characters)": "In the nearest phase, we’ll propose the next project continuing Pilot2-1 from this April-May. Currently, we increase the manufacturing domain-specific use case in the client’s other department and this department, and try to propose the on-premiss version of watsonx.ai to treat more secure data. ",
    "Next steps discussed with the client(limit to 1000 Characters)": "Our client wants to speed up the modeling of the existing system by using this generative AI approach, and finally expand the domain which can utilize the",
    "Is there anything else you want to communicate that will ensure a smooth transition for this client on their CEM journey? ": " ",
}


def append_question(question, user_input_metadata, table_metadata, meta_data_info):
    return v1.prompt_template.format(
        question=question,
        user_input_metadata=user_input_metadata,
        metadata=table_metadata,
        meta_data_info=meta_data_info,
    )


search_item = ["Internal Pre-Sales Engagement", "Close-out and transition details"]
metadata_info = ["internal_pre_sales_data_info", "close_out_transition_info"]


def generate_prompt_answer(
    question_data_dict, data_dict, metadata_info, search, output_dict
):
    prompt = append_question(
        question=question_data_dict,
        user_input_metadata=data_dict,
        table_metadata=input_metadata_question,
        meta_data_info=metadata_info,
    )
    response = model.generate(prompt)
    cleaned_response = clean_response(response, output_dict, search)
    return cleaned_response


def clean_response(generated_response, output_dict, search):
    clean_data = generated_response.get("results")[0].get("generated_text")[:-5]
    smart_data = json.loads(clean_data)
    output_dict[search] = dict()
    for key in smart_data:
        output_dict[search][smart_data.get(key).get("content_key")] = smart_data.get(
            key
        ).get("content_value")
    return output_dict


def extract_data(df):
    data_dict = dict()
    output_dict = dict()
    question_data_dict = dict()
    data = []
    for index, row in df.iterrows():
        data_dict[row["key"]] = row["value"]

    for i, search in enumerate(search_item):
        df_question = stars_data[stars_data["Table"] == search]
        for item in df_question["Information / Question"]:
            question_data_dict[item] = " "
        response = generate_prompt_answer(
            question_data_dict, data_dict, metadata_info[i], search, output_dict
        )
        question_data_dict = dict()
    return response


def structure_response(data):
    table_names = []
    column_categories = []
    column_values = []

    for table_name, content in data.items():
        for category, value in content.items():
            table_names.append(table_name)
            column_categories.append(category)
            column_values.append(value)

    df = pd.DataFrame(
        {
            "Table Name": table_names,
            "Column Category": column_categories,
            "Generated Text ": column_values,
        }
    )
    return df


def extract_table_from_slides(slide):
    i = 0
    data = {}
    for i, shape in enumerate(slide.shapes):
        i += 1
        if shape.has_table:
            table = shape.table
            rows = []
            for i, row in enumerate(table.rows):
                row_data = [cell.text.strip() for cell in row.cells]
                rows.append(row_data)
            df_table = pd.DataFrame(rows, columns=["key", "value"])
            tables = data.get("tables", [])
            tables.append(df_table)
            data["tables"] = tables
        elif hasattr(shape, "text") and not shape.text.isspace() and shape.text != "":
            if len(shape.text_frame.paragraphs) > 1:
                paragraphs = data.get("paragraphs", [])
                for i, paragraph in enumerate(shape.text_frame.paragraphs):
                    if paragraph.text != "":
                        paragraphs.append(paragraph.text)
                data["paragraphs"] = paragraphs
            else:
                text_box = data.get("text_box", [])
                text_box.append(shape.text)
                data["text_box"] = text_box
    data = extract_data(data["tables"][0])

    res = structure_response(data)
    return res


def total_pipeline_process(file_path):
    text_content = ""
    ppt = Presentation(file_path)
    slide = ppt.slides[0]
    res = extract_table_from_slides(slide)
    return res
